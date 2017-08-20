import sys

#import presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from Tkinter import *
from appJar import gui
from os import listdir
from os.path import isfile, join

import time

from find import find

print("application started")

print("define variables")
#define vars
app = gui()

file_name = "Dicsoites: "

mypath= "./songs"
#add files to the list
listin = [f for f in listdir(mypath) if isfile(join(mypath, f))]
listout = []

#define vars end

#functions start here

def listnav(btn):
	print("button \""+btn+"\" pressed")
	if(btn == ">"):
		items= app.getListItems("list_in")
		print(str(items)+" added to output")
		app.addListItems("list_out",items)

	if(btn == "<"):
		items = app.getListItems("list_out")
		for item in items:
			print(str(item)+" removed from output")
			app.removeListItem("list_out",item)

def createprez(btn):
	print("presentation being built")
	songs = []
	#read songs from files
	for path in app.getAllListItems("list_out"):
		songs.extend(open(mypath+"/"+path, "r").readlines());


	prs = Presentation("./template.pptx")
	bullet_slide_layout = prs.slide_layouts[5]

	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes

	body_shape = shapes.placeholders[0]

	tf = body_shape.text_frame
	p = tf.paragraphs[0]

	p.alignment = PP_ALIGN.CENTER


	for song in songs:
		if(song != '\n'):
			print(song)
			p.text += song;
		else:
			bullet_slide_layout = prs.slide_layouts[5]
			slide = prs.slides.add_slide(bullet_slide_layout)
			shapes = slide.shapes
			body_shape = shapes.placeholders[0]

			tf = body_shape.text_frame
			p = tf.paragraphs[0]
			p.alignment = PP_ALIGN.CENTER

	name = file_name+time.strftime("%Y-%m-%d")+".pptx"

	prs.save(name)

	print(name+" mentve")

def search(btn):
	term = app.getEntry('search')
	exact = app.getCheckBox("exact")
	listin = []
	files = [f for f in listdir(mypath) if isfile(join(mypath, f))]

	if(exact):
		print(term)
		for f in files:
			if(find(mypath+"/"+f,term)):
				if f not in listin:
					listin.append(f)

	else:
		terms = term.split()
		print(terms)
		for t in terms:
			for f in files:
				if(find(mypath+"/"+f,t)):
					if f not in listin:
						listin.append(f)

	app.updateListItems("list_in", listin)

def reset_list(btn):
	app.updateListItems("list_in", [f for f in listdir(mypath) if isfile(join(mypath, f))])



#functions end here

# top slice - CREATE the GUI


### fillings go here ###
print("creating ui")

app.addLabel("search", "Search:", 1, 0)
app.addEntry("search", 1, 1)                           # Row 1,Column 1

app.addButton("go",search ,1,2)
app.addButton("reset",reset_list ,2,2)
app.addCheckBox("exact",2,1)

app.addLabel("list", "List:", 3, 0)

app.addListBox("list_in", listin,3,1)

app.addListBox("list_out",listout,3,2)


app.setListBoxMulti("list_in", multi=True)
app.setListBoxMulti("list_out", multi=True)

app.addLabel("fun", "", 3, 0)

app.addButton(">", listnav,4,1)
app.addButton("<", listnav,4,2)
app.addButton("build",createprez,5,0,10)
# bottom slice - START the GUI
print("launching ui")
app.go()
